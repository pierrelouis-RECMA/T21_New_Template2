"""
modern_design.py — Script 1
Design moderne RECMA pour toutes les slides.
Remplace make_native_tables.py.

Palette RECMA :
  Vert foncé header : #2D5C54
  Rouge accent      : #CC2229
  Fond sidebar      : #F4F7F6
  Groupes : Publicis #FFEADD · Omnicom #FFE699 · Dentsu #E2F0D9
            Havas #DCB9FF · WPP #FFE4FF
"""

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Palette ───────────────────────────────────────────────────────────────────
C_HEADER  = '2D5C54'
C_RED     = 'CC2229'
C_WIN     = '1B5E20'
C_DEP     = 'CC2229'
C_WIN_LT  = '2E7D32'
C_SIDEBAR = 'F4F7F6'
C_BORDER  = 'E8EDEB'
C_NBB_POS = '1B5E20'
C_NBB_NEG = 'CC2229'

GROUP_COLORS = {
    'Publicis Media':      'FFEADD',
    'Omnicom Media':       'FFE699',
    'Dentsu':              'E2F0D9',
    'Havas Media Network': 'DCB9FF',
    'WPP Media':           'FFE4FF',
    'IPG Mediabrands':     'D9EAD3', # Ajout pour le Mexique (UM/Initiative)
    'Independent':         'FFFFFF', # Fond Blanc
}
GROUP_BORDER = {
    'Publicis Media':      'FF6B35',
    'Omnicom Media':       'D4A017',
    'Dentsu':              '5A8A3C',
    'Havas Media Network': '7B2FBE',
    'WPP Media':           'B83DB8',
    'IPG Mediabrands':     '38761D',
    'Independent':         'CCCCCC',
}
GROUP_NAMES = {
    'Publicis Media':      'Publicis',
    'Omnicom Media':       'Omnicom',
    'Dentsu':              'Dentsu',
    'Havas Media Network': 'Havas',
    'WPP Media':           'WPP',
    'IPG Mediabrands':     'IPG',
    'Independent':         'Indie',
}

# Couleur de police spécifique pour Independent
C_INDIE_BLUE = '0000FF' # Bleu pur

# ── Helpers ───────────────────────────────────────────────────────────────────
def rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_rect(slide, left, top, width, height, fill_hex=None, line_hex=None, line_w=0):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, left, top, width, height)
    fill  = shape.fill
    if fill_hex:
        fill.solid(); fill.fore_color.rgb = rgb(fill_hex)
    else:
        fill.background()
    line = shape.line
    if line_hex and line_w:
        line.color.rgb = rgb(line_hex); line.width = Pt(line_w)
    else:
        line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             size=9, bold=False, color='1A1A1A', align=PP_ALIGN.LEFT,
             italic=False, wrap=True, fill_hex=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    if fill_hex:
        txBox.fill.solid(); txBox.fill.fore_color.rgb = rgb(fill_hex)
    tf = txBox.text_frame; tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(2)
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    return txBox

def set_cell_fill(cell, hex_color=None):
    tc   = cell._tc
    tcPr = tc.find(qn('a:tcPr'))
    if tcPr is None: tcPr = etree.SubElement(tc, qn('a:tcPr'))
    for tag in ['a:solidFill','a:noFill']:
        for old in tcPr.findall(qn(tag)): tcPr.remove(old)
    if hex_color:
        sf = etree.SubElement(tcPr, qn('a:solidFill'))
        sc = etree.SubElement(sf, qn('a:srgbClr'))
        sc.set('val', hex_color.lstrip('#'))
    else:
        etree.SubElement(tcPr, qn('a:noFill'))

def set_cell_border(cell, color='DDDDDD', width_pt=0.4, sides=None):
    tc   = cell._tc
    tcPr = tc.find(qn('a:tcPr'))
    if tcPr is None: tcPr = etree.SubElement(tc, qn('a:tcPr'))
    if sides is None: sides = ['lnL','lnR','lnT','lnB']
    w = int(width_pt * 12700)
    for side in sides:
        old = tcPr.find(qn(f'a:{side}'))
        if old is not None: tcPr.remove(old)
        ln = etree.SubElement(tcPr, qn(f'a:{side}'))
        ln.set('w', str(w))
        sf = etree.SubElement(ln, qn('a:solidFill'))
        sc = etree.SubElement(sf, qn('a:srgbClr'))
        sc.set('val', color)

def set_cell_left_border(cell, color, width_pt=2.0):
    """Liseré coloré uniquement à gauche de la cellule."""
    tc   = cell._tc
    tcPr = tc.find(qn('a:tcPr'))
    if tcPr is None: tcPr = etree.SubElement(tc, qn('a:tcPr'))
    w = int(width_pt * 12700)
    for side in ['lnL','lnR','lnT','lnB']:
        old = tcPr.find(qn(f'a:{side}'))
        if old is not None: tcPr.remove(old)
        ln   = etree.SubElement(tcPr, qn(f'a:{side}'))
        if side == 'lnL':
            ln.set('w', str(w))
            sf = etree.SubElement(ln, qn('a:solidFill'))
            sc = etree.SubElement(sf, qn('a:srgbClr'))
            sc.set('val', color.lstrip('#'))
        else:
            ln.set('w', str(int(0.3*12700)))
            sf = etree.SubElement(ln, qn('a:solidFill'))
            sc = etree.SubElement(sf, qn('a:srgbClr'))
            sc.set('val', 'E0E0E0')

def write_cell(cell, text, size=8, bold=False, color='1A1A1A',
               align=PP_ALIGN.LEFT, italic=False):
    tf = cell.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(3)
    tf.margin_top  = tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]; p.alignment = align
    for r in p.runs: r.text = ''
    run = p.runs[0] if p.runs else p.add_run()
    run.text = str(text) if text else ''
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)

def get_group_style(group_name):
    bg_color = GROUP_COLORS.get(group_name, 'F2F2F2')
    text_color = C_INDIE_BLUE if group_name == 'Independent' else '000000'
    return bg_color, text_color

def remove_pics(slide):
    spTree = slide.shapes._spTree
    for pic in list(spTree.findall('.//' + qn('p:pic'))):
        spTree.remove(pic)

# ── Header RECMA commun ───────────────────────────────────────────────────────
def add_recma_header(slide, slide_title='', W=None, H=None, market='Hong Kong'):
    """Ajoute le header vert + ligne rouge + motif >> sur toutes les slides."""
    if W is None: W = Emu(7559675)
    if H is None: H = Emu(10439400)

    header_h = Emu(int(H * 0.075))
    red_h    = Emu(int(H * 0.006))

    # Bande verte
    add_rect(slide, 0, 0, W, header_h, fill_hex=C_HEADER)

    # >> rouge
    add_text(slide, '>>', Emu(80000), Emu(20000), Emu(600000), header_h,
             size=16, bold=True, color=C_RED)

    # Titre slide dans header
    full_title = f'{market} Domestic Report  ·  New Business Balance  ·  {slide_title}'
    add_text(slide, full_title,
             Emu(600000), Emu(20000), Emu(W - 1600000), header_h,
             size=8.5, bold=True, color='FFFFFF')

    # Date droite
    add_text(slide, 'March 2026  ·  ©2026 RECMA',
             W - Emu(1500000), Emu(30000), Emu(1400000), header_h,
             size=7, italic=True, color='9FBFBA', align=PP_ALIGN.RIGHT)

    # Ligne rouge
    add_rect(slide, 0, header_h, W, red_h, fill_hex=C_RED)

    return header_h + red_h   # retourne la hauteur totale occupée

# ── Sidebar compacte ──────────────────────────────────────────────────────────
def add_sidebar(slide, top_offset, W, H, label, title, groups=True):
    sidebar_w = Emu(int(W * 0.11))
    content_h = H - top_offset
    inner_w   = sidebar_w - Emu(40000)
    pad       = Emu(80000)

    # Fond sidebar
    add_rect(slide, 0, top_offset, sidebar_w, content_h, fill_hex='F4F7F6')
    # Bordure droite verte
    add_rect(slide, sidebar_w - Emu(30000), top_offset, Emu(30000), content_h,
             fill_hex=C_HEADER)

    y = top_offset + pad

    # Label rouge
    add_text(slide, label, pad//2, y, inner_w, Emu(200000),
             size=7, bold=True, color=C_RED)
    y += Emu(220000)

    # Titre
    add_text(slide, title, pad//2, y, inner_w, Emu(500000),
             size=13, bold=True, color='1A1A1A')
    y += Emu(500000)

    # Ligne rouge séparatrice
    add_rect(slide, pad//2, y, inner_w - pad, Emu(25000), fill_hex=C_RED)
    y += Emu(80000)

    if groups:
        for gname, gcol in GROUP_COLORS.items():
            gborder = GROUP_BORDER[gname]
            gshort  = GROUP_NAMES[gname]
            # Carré couleur
            add_rect(slide, pad//2, y, Emu(120000), Emu(100000),
                     fill_hex=gcol, line_hex=gborder, line_w=1.2)
            # Nom
            add_text(slide, gshort, pad//2 + Emu(140000), y - Emu(10000),
                     inner_w, Emu(130000), size=5.5, color='333333')
            y += Emu(130000)

    # Note bas
    note = 'Retentions & transfers not included · Based on announcement date – 2025'
    add_text(slide, note, pad//2, H - Emu(400000), inner_w, Emu(380000),
             size=4.5, italic=True, color='999999')

    return sidebar_w

# ── Slide 3 — TOP moves (table native) ────────────────────────────────────────
def build_slide3_modern(slide, top_moves, W=None, H=None, market='Hong Kong'):
    if W is None: W = Emu(7559675)
    if H is None: H = Emu(10439400)

    remove_pics(slide)
    top_off = add_recma_header(slide, 'TOP Moves / Retentions 2025', W, H, market=market)
    sidebar_w = add_sidebar(slide, top_off, W, H, 'SLIDE 3', 'TOP\nMoves\n2025', groups=False)

    # Zone table
    tl = sidebar_w + Emu(80000)
    tt = top_off + Emu(120000)
    tw = W - tl - Emu(80000)
    th = H - tt - Emu(150000)

    cols      = ['#','Advertiser','Sector','IS mUSD','Winning Agency','Incumbent','Type']
    col_w_pct = [0.04, 0.22, 0.10, 0.07, 0.20, 0.18, 0.12]
    col_emu   = [Emu(int(tw * p)) for p in col_w_pct]

    n   = len(top_moves) + 1
    tbl = slide.shapes.add_table(n, 7, tl, tt, tw, th).table

    for j, ew in enumerate(col_emu): tbl.columns[j].width = ew
    rh = Emu(int(th / n))
    for i in range(n): tbl.rows[i].height = rh

    # Header
    h_aligns = [PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.LEFT,
                PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.CENTER]
    h_colors = ['FFFFFF','FFFFFF','FFFFFF','A5D6A7','FFFFFF','FFFFFF','FFCDD2']
    for j, (lbl, al, hc) in enumerate(zip(cols, h_aligns, h_colors)):
        cell = tbl.cell(0, j)
        set_cell_fill(cell, C_HEADER)
        set_cell_border(cell, 'FFFFFF', 0.3)
        write_cell(cell, lbl, size=8, bold=True, color=hc, align=al)

    # Données
    for i, (_, row) in enumerate(top_moves.iterrows()):
        ri     = i + 1
        is_ret = row['NewBiz'] == 'RETENTION'
        inc    = str(row.get('Incumbent','') or '')
        vals   = [str(i+1), row['Advertiser'][:28], '',
                  f"{row['Integrated Spends']:.1f}",
                  row['Agency'][:22], inc[:20],
                  'Retention' if is_ret else 'New Business']
        aligns = [PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.LEFT,
                  PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.CENTER]
        colors = [C_HEADER, '1A1A1A', '1A1A1A', C_WIN, C_WIN, '555555',
                  'FF9800' if is_ret else C_HEADER]
        bolds  = [True, is_ret, False, True, True, False, is_ret]

        for j in range(7):
            cell = tbl.cell(ri, j)
            set_cell_fill(cell, None)
            set_cell_border(cell, 'E0E0E0', 0.25)
            write_cell(cell, vals[j], size=7.5, bold=bolds[j],
                       color=colors[j], align=aligns[j])

# ── Slide 4 — T21a agencies (table native) ────────────────────────────────────
def build_slide4_modern(slide, agencies, W=None, H=None, market='Hong Kong'):
    if W is None: W = Emu(7559675)
    if H is None: H = Emu(10439400)

    remove_pics(slide)
    top_off   = add_recma_header(slide, 'T21a — NBB 2025 Agencies Overview', W, H, market=market)
    sidebar_w = add_sidebar(slide, top_off, W, H, 'T21a', 'NBB\n2025')

    tl = sidebar_w + Emu(60000)
    tt = top_off + Emu(100000)
    tw = W - tl - Emu(60000)
    th = H - tt - Emu(130000)

    cols      = ['#','Agency','NBB\n2025 $m','Wins\n$m','Dep.\n$m','W/D',
                 'Main Wins >3$m','Main Dep. >3$m']
    col_w_pct = [0.04, 0.13, 0.07, 0.06, 0.06, 0.05, 0.295, 0.295]
    col_emu   = [Emu(int(tw * p)) for p in col_w_pct]

    n   = len(agencies) + 1
    tbl = slide.shapes.add_table(n, 8, tl, tt, tw, th).table

    for j, ew in enumerate(col_emu): tbl.columns[j].width = ew
    rh = Emu(int(th / n))
    for i in range(n): tbl.rows[i].height = rh

    # Header
    h_al  = [PP_ALIGN.CENTER, PP_ALIGN.LEFT] + [PP_ALIGN.CENTER]*4 + [PP_ALIGN.LEFT]*2
    h_col = ['FFFFFF','FFFFFF','FFFFFF','A5D6A7','FFCDD2','FFFFFF','A5D6A7','FFCDD2']
    for j, (lbl, al, hc) in enumerate(zip(cols, h_al, h_col)):
        cell = tbl.cell(0, j)
        set_cell_fill(cell, C_HEADER)
        set_cell_border(cell, 'FFFFFF', 0.3)
        write_cell(cell, lbl, size=7.5, bold=True, color=hc, align=al)

    # Données
    for i, ag in enumerate(agencies):
        ri   = i + 1
        gbg, g_txt_col = get_group_style(ag['group'])
        gbrd = GROUP_BORDER.get(ag['group'], 'CCCCCC')
        nfg  = C_NBB_POS if ag['nbb'] >= 0 else C_NBB_NEG
        top3 = ri <= 3

        big_w = [r for r in ag['wins_rows'] if r['Integrated Spends'] >= 3]
        big_d = [r for r in ag['dep_rows']  if r['Integrated Spends'] <= -3]
        ws = ' · '.join([f"{r['Advertiser'][:15]} {r['Integrated Spends']:.0f}m" for r in big_w[:3]])
        ds = ' · '.join([f"{r['Advertiser'][:15]} {r['Integrated Spends']:.0f}m" for r in big_d[:3]])

        vals   = [ag['rank'], ag['agency'], f"{ag['nbb']:+.1f}",
                  f"{ag['wins']:.1f}", f"{ag['dep']:.1f}",
                  f"{ag['wc']}/{ag['dc']}", ws, ds]
        aligns = [PP_ALIGN.CENTER, PP_ALIGN.LEFT] + [PP_ALIGN.CENTER]*4 + [PP_ALIGN.LEFT]*2
        colors = [C_RED if top3 else C_HEADER, g_txt_col, nfg,
                  C_WIN, C_DEP, '444444', C_WIN_LT, C_DEP]
        bolds  = [True, True, True, True, True, False, False, False]
        sizes  = [8, 8, 9, 8, 8, 7.5, 6.5, 6.5]

        for j in range(8):
            cell = tbl.cell(ri, j)
            if j == 1:   # Cellule agence : couleur groupe + liseré gauche
                set_cell_fill(cell, gbg)
                set_cell_left_border(cell, gbrd, width_pt=2.5)
            else:
                set_cell_fill(cell, None)
                set_cell_border(cell, 'E0E0E0', 0.25)
            write_cell(cell, vals[j], size=sizes[j], bold=bolds[j],
                       color=colors[j], align=aligns[j])

    # >> filigrane décoratif
    add_text(slide, '>>', W - Emu(900000), H - Emu(1800000), Emu(900000), Emu(1600000),
             size=90, bold=True, color='F0F4F3')
