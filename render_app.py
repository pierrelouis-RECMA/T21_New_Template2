"""
render_app.py — Application Flask pour Render.com
Upload Excel local → génère PPTX + PDF navigable → téléchargement direct

Déploiement sur Render.com :
  1. Push tous les fichiers sur GitHub
  2. Render → New Web Service → connecte ton repo
  3. Build command : pip install -r requirements_render.txt
  4. Start command : gunicorn render_app:app
  5. Partage l'URL à tes collègues

Fichiers requis dans le repo :
  render_app.py
  generate_pptx.py
  modern_design.py
  pdf_nav.py
  T21_HK_pack.pptx
  requirements_render.txt
"""

import os, io, sys, subprocess, shutil, tempfile, logging
from datetime import datetime
from flask import Flask, request, render_template_string, send_file, jsonify

app = Flask(__name__)
logging.basicConfig(level=logging.INFO)
TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), 'T21_HK_pack.pptx')

# ── HTML de l'interface ───────────────────────────────────────────────────────
HTML = """
<!DOCTYPE html>
<html lang="fr">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>NBB Report Generator · RECMA</title>
  <style>
    * { box-sizing: border-box; margin: 0; padding: 0; }
    body { font-family: Arial, sans-serif; background: #F4F7F6; color: #1A1A1A; min-height: 100vh; }

    .header {
      background: #2D5C54; color: white;
      padding: 0 32px; height: 56px;
      display: flex; align-items: center; justify-content: space-between;
      border-bottom: 3px solid #CC2229;
    }
    .header-left { display: flex; align-items: center; gap: 10px; }
    .header-arrow { color: #CC2229; font-size: 22px; font-weight: 900; letter-spacing: -3px; }
    .header-title { font-size: 14px; font-weight: 700; letter-spacing: 1px; }
    .header-badge { background: #CC2229; color: white; font-size: 10px; padding: 2px 8px; border-radius: 3px; }

    .container { max-width: 680px; margin: 48px auto; padding: 0 24px; }

    .card {
      background: white; border-radius: 10px;
      border: 1.5px solid #E0E9E7;
      padding: 32px; margin-bottom: 20px;
    }
    .card-title {
      font-size: 18px; font-weight: 700; color: #2D5C54;
      margin-bottom: 6px;
    }
    .card-sub { font-size: 13px; color: #666; margin-bottom: 24px; line-height: 1.6; }

    .red-line { height: 2px; background: #CC2229; margin: 16px 0; }

    label { font-size: 13px; font-weight: 700; color: #2D5C54; display: block; margin-bottom: 6px; }

    .upload-zone {
      border: 2px dashed #9FBFBA; border-radius: 8px;
      padding: 32px; text-align: center; cursor: pointer;
      background: #F9FBFB; transition: all 0.2s;
      margin-bottom: 20px;
    }
    .upload-zone:hover { border-color: #2D5C54; background: #F0F6F5; }
    .upload-zone.active { border-color: #CC2229; background: #FFF5F5; }
    .upload-zone input { display: none; }
    .upload-icon { font-size: 32px; margin-bottom: 8px; }
    .upload-text { font-size: 14px; color: #555; }
    .upload-sub  { font-size: 12px; color: #999; margin-top: 4px; }
    .file-name   { font-size: 13px; font-weight: 700; color: #2D5C54; margin-top: 8px; }

    .options { display: grid; grid-template-columns: 1fr 1fr; gap: 12px; margin-bottom: 20px; }
    select, input[type=text] {
      width: 100%; padding: 9px 12px; border: 1.5px solid #E0E9E7;
      border-radius: 6px; font-size: 13px; color: #1A1A1A; background: white;
    }
    select:focus, input:focus { outline: none; border-color: #2D5C54; }

    .checkboxes { display: flex; gap: 20px; margin-bottom: 24px; }
    .check-item { display: flex; align-items: center; gap: 6px; font-size: 13px; cursor: pointer; }
    .check-item input { width: 16px; height: 16px; accent-color: #2D5C54; }

    .btn {
      display: block; width: 100%; padding: 14px;
      background: #2D5C54; color: white; border: none; border-radius: 8px;
      font-size: 15px; font-weight: 700; cursor: pointer; letter-spacing: 0.5px;
      transition: background 0.2s;
    }
    .btn:hover { background: #1E4039; }
    .btn:disabled { background: #9FBFBA; cursor: not-allowed; }
    .btn-red { background: #CC2229; }
    .btn-red:hover { background: #A31B21; }

    .progress {
      display: none; background: #F0F6F5; border-radius: 8px;
      padding: 16px; margin-top: 16px; text-align: center;
    }
    .spinner {
      display: inline-block; width: 20px; height: 20px;
      border: 3px solid #E0E9E7; border-top-color: #2D5C54;
      border-radius: 50%; animation: spin 0.8s linear infinite;
      margin-right: 8px; vertical-align: middle;
    }
    @keyframes spin { to { transform: rotate(360deg); } }

    .downloads { display: none; margin-top: 16px; }
    .dl-btn {
      display: flex; align-items: center; gap: 10px;
      padding: 12px 16px; border-radius: 8px; margin-bottom: 10px;
      text-decoration: none; font-weight: 700; font-size: 14px;
      transition: opacity 0.2s;
    }
    .dl-btn:hover { opacity: 0.85; }
    .dl-pptx { background: #CC2229; color: white; }
    .dl-pdf  { background: #2D5C54; color: white; }
    .dl-icon { font-size: 20px; }

    .info-box {
      background: #F0F6F5; border-left: 3px solid #2D5C54;
      border-radius: 0 6px 6px 0; padding: 12px 14px;
      font-size: 12px; color: #555; line-height: 1.7;
    }
    .info-box strong { color: #2D5C54; }

    .error-box {
      background: #FFF0F0; border-left: 3px solid #CC2229;
      border-radius: 0 6px 6px 0; padding: 12px 14px;
      font-size: 13px; color: #CC2229; display: none; margin-top: 12px;
    }

    footer { text-align: center; color: #999; font-size: 11px; margin: 32px 0; }
  </style>
</head>
<body>

<div class="header">
  <div class="header-left">
    <span class="header-arrow">>></span>
    <span class="header-title">NBB REPORT GENERATOR</span>
  </div>
  <span class="header-badge">RECMA · Hong Kong</span>
</div>

<div class="container">

  <div class="card">
    <div class="card-title">Générer le rapport New Business Balance</div>
    <div class="card-sub">
      Upload ton fichier Excel source → reçois le PPTX + PDF en 30 secondes.<br>
      Aucune donnée n'est conservée sur le serveur.
    </div>
    <div class="red-line"></div>

    <!-- Upload -->
    <label>Fichier Excel source</label>
    <div class="upload-zone" id="dropZone" onclick="document.getElementById('fileInput').click()">
      <div class="upload-icon">📊</div>
      <div class="upload-text">Glisse ton fichier ici ou clique pour parcourir</div>
      <div class="upload-sub">Formats acceptés : .xlsx, .xls</div>
      <div class="file-name" id="fileName"></div>
      <input type="file" id="fileInput" accept=".xlsx,.xls" onchange="onFileSelect(this)">
    </div>

    <!-- Options -->
    <div class="options">
      <div>
        <label>Marché</label>
        <select id="market">
          <option value="Hong_Kong">Hong Kong</option>
          <option value="Indonesia">Indonesia</option>
          <option value="Singapore">Singapore</option>
          <option value="Vietnam">Vietnam</option>
          <option value="Vietnam">Mexico</option>
          <option value="Malaysia">Malaysia</option>
        </select>
      </div>
      <div>
        <label>Année</label>
        <input type="text" id="year" value="2025">
      </div>
    </div>

    <!-- Outputs -->
    <label>Fichiers à générer</label>
    <div class="checkboxes">
      <label class="check-item">
        <input type="checkbox" id="genPptx" checked> PPTX (éditable)
      </label>
      <label class="check-item">
        <input type="checkbox" id="genPdf"  checked> PDF (navigation cliquable)
      </label>
    </div>

    <!-- Bouton -->
    <button class="btn" id="genBtn" onclick="generate()" disabled>
      Générer le rapport
    </button>

    <!-- Progress -->
    <div class="progress" id="progress">
      <span class="spinner"></span>
      <span id="progressText">Génération en cours…</span>
    </div>

    <!-- Pays détecté -->
    <div id="marketDetectedBox" style="display:none; margin-top:15px; padding:10px; background:#E8F5E9; border-left:4px solid #4CAF50; font-size:14px;">
      🌍 Pays détecté : <strong id="detectedMarketName"></strong>
    </div>

    <!-- Erreur -->
    <div class="error-box" id="errorBox"></div>

    <!-- Téléchargements -->
    <div class="downloads" id="downloads">
      <a class="dl-btn dl-pptx" id="dlPptx" href="#" download>
        <span class="dl-icon">📑</span> Télécharger le PPTX
      </a>
      <a class="dl-btn dl-pdf" id="dlPdf" href="#" download>
        <span class="dl-icon">📄</span> Télécharger le PDF navigable
      </a>
    </div>
  </div>

  <!-- Info -->
  <div class="card">
    <div class="info-box">
      <strong>Colonnes Excel attendues :</strong><br>
      Agency · NewBiz (WIN/DEPARTURE/RETENTION) · Advertiser ·
      Integrated Spends · Ad Spends · Date of announcement ·
      Country of Decision · Incumbent
      <br><br>
      <strong>PDF navigable :</strong> Table des matières cliquable · Bookmarks dans la sidebar ·
      Boutons Précédent/Suivant en bas de page
    </div>
  </div>

</div>

<footer>©2026 RECMA · NBB Report Generator · Confidentiel</footer>

<script>
  let selectedFile = null;
  let pptxBlob = null, pdfBlob = null;

  // Drag & drop
  const dz = document.getElementById('dropZone');
  dz.addEventListener('dragover',  e => { e.preventDefault(); dz.classList.add('active'); });
  dz.addEventListener('dragleave', () => dz.classList.remove('active'));
  dz.addEventListener('drop', e => {
    e.preventDefault(); dz.classList.remove('active');
    const f = e.dataTransfer.files[0];
    if (f) setFile(f);
  });

  function onFileSelect(input) {
    if (input.files[0]) setFile(input.files[0]);
  }
  function setFile(f) {
    selectedFile = f;
    document.getElementById('fileName').textContent = '✓ ' + f.name;
    document.getElementById('genBtn').disabled = false;
  }

  async function generate() {
    if (!selectedFile) return;

    const btn       = document.getElementById('genBtn');
    const progress  = document.getElementById('progress');
    const downloads = document.getElementById('downloads');
    const errorBox  = document.getElementById('errorBox');
    const progText  = document.getElementById('progressText');
    const marketBox = document.getElementById('marketDetectedBox');

    btn.disabled = true;
    progress.style.display  = 'block';
    downloads.style.display = 'none';
    errorBox.style.display  = 'none';

    const steps = [
      'Chargement des données Excel…',
      'Calcul des statistiques NBB…',
      'Génération des tableaux PPTX…',
      'Finalisation du rapport…',
    ];
    let step = 0;
    const interval = setInterval(() => {
      progText.textContent = steps[step % steps.length];
      step++;
    }, 4000);

    try {
      const fd = new FormData();
      fd.append('file',   selectedFile);
      fd.append('market', document.getElementById('market').value);
      fd.append('year',   document.getElementById('year').value);
      fd.append('pptx',   document.getElementById('genPptx').checked ? '1' : '0');
      fd.append('pdf',    document.getElementById('genPdf').checked  ? '1' : '0');

      const resp = await fetch('/generate', { method: 'POST', body: fd });

      if (!resp.ok) {
        const err = await resp.json();
        throw new Error(err.error || 'Erreur serveur');
      }

      const data = await resp.json();

      // Affichage du pays détecté
      marketBox.style.display = 'block';
      document.getElementById('detectedMarketName').textContent = data.market_detected;

      // Téléchargements
      downloads.style.display = 'block';
      const dlPptx = document.getElementById('dlPptx');
      const dlPdf  = document.getElementById('dlPdf');

      if (data.pptx_url) {
        dlPptx.href = data.pptx_url;
        dlPptx.download = data.pptx_filename;
        dlPptx.style.display = 'flex';
      } else {
        dlPptx.style.display = 'none';
      }

      if (data.pdf_url) {
        dlPdf.href = data.pdf_url;
        dlPdf.download = data.pdf_filename;
        dlPdf.style.display = 'flex';
      } else {
        dlPdf.style.display = 'none';
      }

    } catch (e) {
      errorBox.style.display = 'block';
      errorBox.textContent   = '❌ ' + e.message;
    } finally {
      clearInterval(interval);
      progress.style.display = 'none';
      btn.disabled = false;
    }
  }
</script>
</body>
</html>
"""

# ── Routes ────────────────────────────────────────────────────────────────────
@app.route('/')
def index():
    return render_template_string(HTML)

@app.route('/generate', methods=['POST'])
def generate():
    import uuid, os
    from pptx import Presentation

    try:
        if 'file' not in request.files:
            return jsonify({'error': 'Aucun fichier uploadé'}), 400

        f       = request.files['file']
        market  = request.form.get('market', 'Hong_Kong')
        year    = request.form.get('year',   '2025')
        do_pptx = request.form.get('pptx',  '1') == '1'
        do_pdf  = request.form.get('pdf',   '1') == '1'

        if not f.filename.endswith(('.xlsx', '.xls')):
            return jsonify({'error': 'Format non supporté. Utilise .xlsx ou .xls'}), 400

        with tempfile.TemporaryDirectory() as tmp:
            excel_path = os.path.join(tmp, 'data.xlsx')
            f.save(excel_path)

            # Charge le module generate_pptx
            sys.path.insert(0, os.path.dirname(__file__))
            import generate_pptx as gp

            # Passe le chemin du fichier uploadé directement
            gp.EXCEL = excel_path

            # Stats
            agencies, group_stats, top_moves, df, market_detected = gp.load_stats()
            # Priorité au marché détecté dans l'Excel s'il est cohérent
            market_to_use = market_detected if market_detected != "Hong Kong" else market.replace('_', ' ')
            
            app.logger.info(f"{len(agencies)} agences, {len(df)} lignes, pays: {market_to_use}")

            # Images slides 5/6/7
            imgs = {
                3: gp.make_slide3_img(top_moves),
                4: gp.make_slide4_img(agencies),
                5: gp.make_slide5_img(agencies, group_stats),
                6: gp.make_slide6_chart(agencies),
                7: gp.make_slide7_img(agencies),
            }

            # PPTX
            prs = Presentation(TEMPLATE_PATH)
            gp.replace_text_in_pptx(prs, "Hong Kong", market_to_use)

            gp.update_slide2(prs, agencies, group_stats)

            # Images slides 3-7 (version originale)
            for idx in [2, 3, 4, 5, 6]:
                replaced = gp.replace_slide_image(prs, idx, imgs[idx+1], None)
                if not replaced and idx == 5:
                    slide = prs.slides[5]; imgs[6].seek(0)
                    for shape in slide.shapes:
                        if shape.shape_type == 3:
                            slide.shapes.add_picture(imgs[6], shape.left, shape.top,
                                                      shape.width, shape.height); break

            ts         = datetime.now().strftime('%Y%m%d_%H%M')
            base_name  = f"NBB_{market_to_use.replace(' ', '_')}_{year}_{ts}"
            result     = {}

            if do_pptx:
                pptx_path = os.path.join(tmp, f"{base_name}.pptx")
                prs.save(pptx_path)

                # Stocke pour téléchargement
                token = str(uuid.uuid4())
                with open(pptx_path, 'rb') as fp:
                    _file_cache[token] = (fp.read(), f"{base_name}.pptx",
                                          'application/vnd.openxmlformats-officedocument.presentationml.presentation')

                result['pptx_url']      = f'/download/{token}'
                result['pptx_filename'] = f"{base_name}.pptx"

                if do_pdf:
                    # Converti en PDF avec LibreOffice
                    pdf_raw = os.path.join(tmp, f"{base_name}.pdf")
                    subprocess.run(
                        ['libreoffice', '--headless',
                         '--convert-to', 'pdf',
                         '--outdir', tmp, pptx_path],
                        capture_output=True, timeout=180
                    )
                    # LibreOffice génère le PDF avec le même nom de base
                    pdf_gen = pptx_path.replace('.pptx', '.pdf')

                    if os.path.exists(pdf_gen):
                        # Ajoute la navigation PDF
                        pdf_nav_path = os.path.join(tmp, f"{base_name}_nav.pdf")
                        from pdf_nav import add_pdf_navigation
                        add_pdf_navigation(pdf_gen, pdf_nav_path)

                        token_pdf = str(uuid.uuid4())
                        with open(pdf_nav_path, 'rb') as fp:
                            _file_cache[token_pdf] = (fp.read(), f"{base_name}.pdf", 'application/pdf')

                        result['pdf_url']      = f'/download/{token_pdf}'
                        result['pdf_filename'] = f"{base_name}.pdf"

            result['market_detected'] = market_to_use

        return jsonify(result)

    except Exception as e:
        app.logger.error(f"Erreur: {e}", exc_info=True)
        return jsonify({'error': str(e)}), 500

# ── Cache mémoire (tokens) ────────────────────────────────────────────────────
_file_cache = {}  # token → (bytes, filename, mimetype)

@app.route('/download/<token>')
def download(token):
    if token not in _file_cache:
        return 'Lien expiré ou invalide', 404
    data, filename, mime = _file_cache[token]
    return send_file(io.BytesIO(data), mimetype=mime,
                     as_attachment=True, download_name=filename)

@app.route('/health')
def health():
    return jsonify({'status': 'ok', 'template': os.path.exists(TEMPLATE_PATH)})

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=int(os.environ.get('PORT', 5000)), debug=False)
