"""
NBB Report – PPTX Generator
Prend le template T21_HK_pack.pptx et le remplit avec les données du fichier Excel.
Usage : python3 generate_pptx.py
"""

import copy, io, os, shutil, warnings
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib import rcParams
rcParams['font.family'] = 'DejaVu Sans'

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

warnings.filterwarnings('ignore')

# ── Paths ────────────────────────────────────────────────────────────────────
TEMPLATE  = '/mnt/user-data/uploads/T21_HK_Agencies_Glass_v12.pptx'
EXCEL     = '/mnt/user-data/uploads/Newbiz_Balance_DB_Report_HK_2025.xlsx'
OUTPUT    = '/mnt/user-data/outputs/HK_NBB_Report_2025_AUTO.pptx'
IMG_DIR   = '/home/claude/slide_imgs'
os.makedirs(IMG_DIR, exist_ok=True)

# ── Palette ──────────────────────────────────────────────────────────────────
C_HEADER   = '#2D5C54'
C_WIN      = '#4CAF50'
C_DEP      = '#E53935'
C_RET      = '#FF9800'
C_ALT      = '#F2F6F5'
C_TOTAL    = '#E8F5E9'
C_SUBHDR   = '#3D7068'
C_NBB_POS  = '#1B5E20'
C_NBB_NEG  = '#B71C1C'

GROUP_COLORS = {
    'Publicis Media':      '#FFEADD',
    'Omnicom Media':       '#FFE699',
    'Dentsu':              '#E2F0D9',
    'Havas Media Network': '#DCB9FF',
    'WPP Media':           '#FFE4FF',
    'IPG Mediabrands':     '#D9EAD3',
    'Independent':         '#FFFFFF',
}
GROUP_ORDER = ['Publicis Media','Omnicom Media','Dentsu','Havas Media Network','WPP Media', 'IPG Mediabrands', 'Independent']

AGENCY_GROUP = {
    # Agences existantes
    'SPARK FOUNDRY': 'Publicis Media', 'STARCOM': 'Publicis Media', 'ZENITH': 'Publicis Media',
    'PHD': 'Omnicom Media', 'OMD': 'Omnicom Media', 'HEARTS & SCIENCE': 'Omnicom Media',
    'DENTSU X': 'Dentsu', 'IPROSPECT': 'Dentsu', 'CARAT': 'Dentsu',
    'HAVAS MEDIA': 'Havas Media Network',
    'ESSENCEMEDIACOM': 'WPP Media', 'WAVEMAKER': 'WPP Media', 'MINDSHARE': 'WPP Media',
    
    # Nouvelles agences Mexique / Global
    'ARENA': 'Havas Media Network',
    'VALE MEDIA': 'Independent',
    'INITIATIVE': 'IPG Mediabrands',
    'UM': 'IPG Mediabrands',
}

def get_agency_group(agency_name):
    """Nettoie le nom de l'agence et retourne son groupe (ou Independent par défaut)."""
    clean_name = str(agency_name).upper().strip()
    return AGENCY_GROUP.get(clean_name, 'Independent')

# ── Load & process data ───────────────────────────────────────────────────────
def load_stats():
    # Lit le premier onglet quel que soit son nom
    xl   = pd.ExcelFile(EXCEL)
    sheet = xl.sheet_names[0]
    df   = pd.read_excel(EXCEL, sheet_name=sheet)
    df['Agency'] = df['Agency'].str.strip().str.upper()
    df['NewBiz']  = df['NewBiz'].str.strip().str.upper()
    df['Group']   = df['Agency'].apply(get_agency_group)

    # Détection automatique du Pays
    country_col = 'Country' if 'Country' in df.columns else 'Country of Decision'
    countries = df[country_col].dropna().unique() if country_col in df.columns else []
    market_name = str(countries[0]) if len(countries) > 0 else "Hong Kong"

    def fmt_date(v):
        try:
            if isinstance(v, (int, float)):
                d = pd.to_datetime(v, origin='1899-12-30', unit='D')
            else:
                d = pd.to_datetime(v)
            return d.strftime('%b-%y')
        except: return ''

    df['Date_str'] = df['Date of announcement'].apply(fmt_date)

    wins_df = df[df.NewBiz=='WIN']
    dep_df  = df[df.NewBiz=='DEPARTURE']
    ret_df  = df[df.NewBiz=='RETENTION']

    agencies = []
    for ag in df.Agency.unique():
        w = wins_df[wins_df.Agency==ag]['Integrated Spends'].sum()
        d = dep_df [dep_df.Agency==ag]['Integrated Spends'].sum()
        r = ret_df [ret_df.Agency==ag]['Integrated Spends'].sum()
        agencies.append({
            'agency': ag, 'group': get_agency_group(ag),
            'nbb': w+d, 'wins': w, 'dep': d, 'ret': r,
            'wc': len(wins_df[wins_df.Agency==ag]),
            'dc': len(dep_df [dep_df.Agency==ag]),
            'rc': len(ret_df [ret_df.Agency==ag]),
            'wins_rows': wins_df[wins_df.Agency==ag].to_dict('records'),
            'dep_rows':  dep_df [dep_df.Agency==ag].to_dict('records'),
            'ret_rows':  ret_df [ret_df.Agency==ag].to_dict('records'),
        })
    agencies.sort(key=lambda x: -x['nbb'])
    for i,a in enumerate(agencies): a['rank'] = i+1

    group_stats = {}
    for g in GROUP_ORDER:
        ags = [a for a in agencies if a['group']==g]
        group_stats[g] = {
            'group':g, 'agencies':ags,
            'nbb':  sum(a['nbb']  for a in ags),
            'wins': sum(a['wins'] for a in ags),
            'dep':  sum(a['dep']  for a in ags),
            'wc':   sum(a['wc']   for a in ags),
            'dc':   sum(a['dc']   for a in ags),
        }

    top_moves = df[df.NewBiz.isin(['WIN','RETENTION'])].copy()
    top_moves['IS_abs'] = top_moves['Integrated Spends'].abs()
    top_moves = top_moves.sort_values('IS_abs', ascending=False).head(24)

    return agencies, group_stats, top_moves, df, market_name

def r2(v): return round(v, 1)

def replace_text_in_pptx(prs, old_text, new_text):
    """Remplace dynamiquement le texte dans tout le template."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if old_text in run.text:
                            run.text = run.text.replace(old_text, new_text)

# ── Matplotlib table helpers ──────────────────────────────────────────────────
def hex2rgb(h): h=h.lstrip('#'); return tuple(int(h[i:i+2],16)/255 for i in (0,2,4))

def make_table_img(col_labels, col_widths, rows, row_styles,
                   title='', subtitle='', figw=10.5, dpi=150):
    """
    rows: list of lists (cell values)
    row_styles: list of dicts with keys: bg, fg, bold, size (per row)
    """
    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(col_labels)
    row_h  = 0.22
    title_h = 0.35 if title else 0
    sub_h   = 0.22 if subtitle else 0
    fig_h   = title_h + sub_h + n_rows * row_h + 0.15

    fig, ax = plt.subplots(figsize=(figw, fig_h))
    fig.patch.set_facecolor('white')
    ax.set_axis_off()
    ax.set_xlim(0, figw)
    ax.set_ylim(0, fig_h)

    y_offset = fig_h

    # Title
    if title:
        y_offset -= title_h * 0.5
        ax.text(0.05, y_offset, title, fontsize=11, fontweight='bold',
                color=C_HEADER, va='center')
        y_offset -= title_h * 0.5

    # Subtitle
    if subtitle:
        y_offset -= sub_h * 0.5
        ax.text(0.05, y_offset, subtitle, fontsize=7.5, color='#555555',
                va='center', style='italic')
        y_offset -= sub_h * 0.5

    # Compute x positions from col widths (normalized to figw)
    total_w = sum(col_widths)
    x_starts = []
    x = 0.0
    for w in col_widths:
        x_starts.append(x / total_w * figw)
        x += w
    x_starts.append(figw)  # right edge

    def draw_row(y_top, height, cells, style, is_header=False):
        bg = style.get('bg', 'white')
        fg = style.get('fg', 'black')
        bold= style.get('bold', False)
        sz  = style.get('size', 8)
        aligns = style.get('aligns', ['left']*n_cols)
        fg_per = style.get('fg_per', None)  # per-cell fg override

        # Draw background
        rect = mpatches.FancyBboxPatch(
            (0, y_top - height), figw, height,
            boxstyle='square,pad=0', linewidth=0,
            facecolor=hex2rgb(bg) if bg!='white' else 'white', zorder=1)
        ax.add_patch(rect)

        # Border bottom
        ax.plot([0, figw], [y_top - height, y_top - height],
                color='#CCCCCC', linewidth=0.4, zorder=2)

        # Cells
        for j, (cell, x0, x1) in enumerate(zip(cells, x_starts[:-1], x_starts[1:])):
            al  = aligns[j] if j < len(aligns) else 'left'
            cell_fg = fg
            if fg_per and j < len(fg_per) and fg_per[j]: cell_fg = fg_per[j]

            xp = x0 + 0.05 if al == 'left' else (x0+x1)/2 if al == 'center' else x1 - 0.05
            ha = 'left' if al == 'left' else 'center' if al == 'center' else 'right'
            yp = y_top - height / 2

            ax.text(xp, yp, str(cell), fontsize=sz,
                    color=cell_fg, ha=ha, va='center',
                    fontweight='bold' if bold else 'normal',
                    zorder=3, clip_on=True)

        # Vertical dividers
        for xd in x_starts[1:-1]:
            ax.plot([xd, xd], [y_top - height, y_top],
                    color='#DDDDDD', linewidth=0.3, zorder=2)

    # Header row
    y = y_offset
    draw_row(y, row_h,
             col_labels,
             {'bg': C_HEADER, 'fg': 'white', 'bold': True, 'size': 8,
              'aligns': ['left'] + ['center']*(n_cols-1)},
             is_header=True)
    y -= row_h

    # Data rows
    for i, (row_vals, style) in enumerate(zip(rows, row_styles)):
        if 'bg' not in style:
            style['bg'] = 'white' if i % 2 == 0 else C_ALT
        draw_row(y, row_h, row_vals, style)
        y -= row_h

    # Top border
    ax.plot([0, figw], [y_offset, y_offset], color=C_HEADER, linewidth=1.2, zorder=4)
    # Outer border
    ax.plot([0,figw,figw,0,0],[y_offset,y_offset,
             y_offset-n_rows*row_h,y_offset-n_rows*row_h,y_offset],
            color='#AAAAAA', linewidth=0.5, zorder=4)

    plt.tight_layout(pad=0)
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close()
    buf.seek(0)
    return buf

# ── Slide 3 — TOP moves ───────────────────────────────────────────────────────
def make_slide3_img(top_moves):
    cols    = ['#','Advertiser','Sector','IS mUSD','Winning Agency','Incumbent','Type']
    widths  = [0.5, 2.8, 1.5, 1.0, 1.8, 1.5, 1.3]
    rows, styles = [], []

    # Sector mapping from incumbent/assignment col
    for i, (_, row) in enumerate(top_moves.iterrows()):
        is_ret = row['NewBiz'] == 'RETENTION'
        inc_df = top_moves[(top_moves['Advertiser']==row['Advertiser'])
                           & (top_moves['NewBiz']=='DEPARTURE')]
        inc = row.get('Incumbent','') or ''

        r = [i+1,
             row['Advertiser'][:30],
             '',
             f"{row['Integrated Spends']:.1f}",
             row['Agency'][:20],
             str(inc)[:18],
             'Retention' if is_ret else 'New Business']

        fg_type = C_RET if is_ret else C_WIN
        rows.append(r)
        styles.append({
            'bg': '#FFF8EE' if is_ret else ('white' if i%2==0 else C_ALT),
            'size': 7.5,
            'aligns': ['center','left','left','center','left','left','center'],
            'fg_per': [None, C_RET if is_ret else None, None,
                       C_WIN, C_WIN, None, fg_type],
            'bold': is_ret,
        })

    return make_table_img(cols, widths, rows, styles,
                          title='TOP moves / retentions 2025',
                          subtitle='Based on Integrated Spendings – 2025',
                          figw=10.8, dpi=160)

# ── Slide 4 — T21a agencies overview ─────────────────────────────────────────
def make_slide4_img(agencies):
    cols   = ['Rkg','Agency','NBB\n2025 $m','Wins\n$m','Dep.\n$m','W/D','Main Wins >3$m','Main Dep. >3$m']
    widths = [0.45, 1.55, 0.75, 0.65, 0.65, 0.5, 2.8, 2.8]
    rows, styles = [], []

    for ag in agencies:
        big_w = [r for r in ag['wins_rows'] if r['Integrated Spends']>=3]
        big_d = [r for r in ag['dep_rows']  if r['Integrated Spends']<=-3]
        ws = ' | '.join([f"{r['Advertiser'][:14]} {r['Integrated Spends']:.0f}m" for r in big_w[:3]])
        ds = ' | '.join([f"{r['Advertiser'][:14]} {r['Integrated Spends']:.0f}m" for r in big_d[:3]])
        gbg = GROUP_COLORS.get(ag['group'], 'white')
        nbb_fg = C_NBB_POS if ag['nbb']>=0 else C_NBB_NEG

        rows.append([ag['rank'], ag['agency'][:18],
                     r2(ag['nbb']), r2(ag['wins']), r2(ag['dep']),
                     f"{ag['wc']}/{ag['dc']}", ws, ds])
        styles.append({
            'bg': gbg, 'size': 7,
            'aligns': ['center','left','center','center','center','center','left','left'],
            'fg_per': [None, None, nbb_fg, C_WIN, C_DEP, None, C_WIN, C_DEP],
        })

    return make_table_img(cols, widths, rows, styles,
                          title='T21a – NBB 2025 agencies overview',
                          subtitle='Retentions, contract renewals & transfers not included | Based on date of announcement – 2025',
                          figw=10.8, dpi=160)

# ── Slide 5 — T21b group overview ────────────────────────────────────────────
def make_slide5_img(agencies, group_stats):
    cols   = ['Rkg','Group / Agency','NBB 2025 $m','Wins $m','Dep. $m','W / D']
    widths = [0.45, 2.8, 1.2, 1.0, 1.0, 0.8]
    rows, styles = [], []

    for rank, gname in enumerate(GROUP_ORDER, 1):
        gs   = group_stats[gname]
        gbg  = GROUP_COLORS.get(gname,'#EEEEEE')
        nfg  = C_NBB_POS if gs['nbb']>=0 else C_NBB_NEG
        rows.append([rank, gname, r2(gs['nbb']), r2(gs['wins']), r2(gs['dep']),
                     f"{gs['wc']}/{gs['dc']}"])
        styles.append({'bg':gbg,'bold':True,'size':8.5,
                       'aligns':['center','left','center','center','center','center'],
                       'fg_per':[None,None,nfg,C_WIN,C_DEP,None]})
        for ag in gs['agencies']:
            nfg2 = C_NBB_POS if ag['nbb']>=0 else C_NBB_NEG
            rows.append(['', f"   {ag['agency']}", r2(ag['nbb']), r2(ag['wins']),
                         r2(ag['dep']), f"{ag['wc']}/{ag['dc']}"])
            styles.append({'bg':'white','size':7.5,'bold':False,
                           'aligns':['center','left','center','center','center','center'],
                           'fg_per':[None,None,nfg2,C_WIN,C_DEP,None]})

    # Total
    tN=sum(group_stats[g]['nbb']  for g in GROUP_ORDER)
    tW=sum(group_stats[g]['wins'] for g in GROUP_ORDER)
    tD=sum(group_stats[g]['dep']  for g in GROUP_ORDER)
    tWC=sum(group_stats[g]['wc']  for g in GROUP_ORDER)
    tDC=sum(group_stats[g]['dc']  for g in GROUP_ORDER)
    rows.append(['','TOTAL 5 groups',r2(tN),r2(tW),r2(tD),f"{tWC}/{tDC}"])
    styles.append({'bg':C_TOTAL,'bold':True,'size':9,
                   'aligns':['center','left','center','center','center','center'],
                   'fg_per':[None,None,C_NBB_POS if tN>=0 else C_NBB_NEG,C_WIN,C_DEP,None]})

    return make_table_img(cols, widths, rows, styles,
                          title='T21b – NBB 2025 Group overview',
                          subtitle='Retentions, contract renewals & transfers not included | Based on date of announcement – 2025',
                          figw=7.5, dpi=160)

# ── Slide 6 — Retentions chart ────────────────────────────────────────────────
def make_slide6_chart(agencies):
    ret_data = [(a['agency'], a['ret']) for a in agencies if a['ret'] != 0]
    ret_data.sort(key=lambda x: x[1], reverse=True)

    fig, ax = plt.subplots(figsize=(7.5, 3.2))
    fig.patch.set_facecolor('white')
    ax.set_facecolor('white')

    ags  = [r[0] for r in ret_data]
    vals = [r[1] for r in ret_data]
    bar_colors = ['#FFECB3' if v>=0 else '#FFCDD2' for v in vals]

    bars = ax.barh(range(len(ags)), vals, color=bar_colors, edgecolor='white', height=0.6)
    ax.set_yticks(range(len(ags)))
    ax.set_yticklabels(ags, fontsize=9)
    ax.invert_yaxis()
    ax.axvline(0, color='#999', linewidth=0.8)

    for i, (bar, val) in enumerate(zip(bars, vals)):
        xp = val + 0.3 if val>=0 else val-0.3
        ha = 'left' if val>=0 else 'right'
        ax.text(xp, i, f'{val:.0f}', va='center', ha=ha,
                fontsize=9, fontweight='bold', color='#333')

    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_visible(False)
    ax.set_xlabel('Retentions $m', fontsize=8.5)
    ax.set_title('T21c – NBB 2025 Retentions ranking by agency',
                 fontsize=11, color=C_HEADER, fontweight='bold', pad=8)
    plt.tight_layout(pad=0.5)

    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor='white')
    plt.close()
    buf.seek(0)
    return buf

# ── Slide 7 — T21d details ────────────────────────────────────────────────────
def make_slide7_img(agencies):
    cols   = ['Agency','NBB\n$m','WINS — Advertiser','IS\n$m','Date','DEPARTURES — Advertiser','IS\n$m']
    widths = [1.4, 0.7, 2.5, 0.6, 0.65, 2.5, 0.6]
    rows, styles = [], []

    for ag in agencies:
        gbg = GROUP_COLORS.get(ag['group'],'#F5F5F5')
        max_r = max(len(ag['wins_rows']), len(ag['dep_rows']), 1)

        for i in range(max_r):
            wr = ag['wins_rows'][i] if i < len(ag['wins_rows']) else None
            dr = ag['dep_rows'][i]  if i < len(ag['dep_rows'])  else None
            nfg = C_NBB_POS if ag['nbb']>=0 else C_NBB_NEG

            rows.append([
                ag['agency'][:16] if i==0 else '',
                r2(ag['nbb']) if i==0 else '',
                wr['Advertiser'][:26] if wr else '',
                r2(wr['Integrated Spends']) if wr else '',
                wr['Date_str'] if wr else '',
                dr['Advertiser'][:26] if dr else '',
                r2(dr['Integrated Spends']) if dr else '',
            ])
            styles.append({
                'bg': gbg if i==0 else ('white' if i%2==0 else C_ALT),
                'bold': i==0, 'size': 7,
                'aligns':['left','center','left','center','center','left','center'],
                'fg_per':[None, nfg if i==0 else None,
                          C_WIN if wr else None, C_WIN if wr else None, None,
                          C_DEP if dr else None, C_DEP if dr else None],
            })

        # Total row
        rows.append([f'Total {ag["agency"][:12]}','',
                     f'{ag["wc"]} win(s)', r2(ag['wins']),'',
                     f'{ag["dc"]} dep.', r2(ag['dep'])])
        styles.append({'bg':C_TOTAL,'bold':True,'size':7,
                       'aligns':['left','center','left','center','center','left','center'],
                       'fg_per':[None,None,C_WIN,C_WIN,None,C_DEP,C_DEP]})

    return make_table_img(cols, widths, rows, styles,
                          title='T21d – NBB 2025 details by agency',
                          subtitle='Retentions, contract renewals & transfers not included | Based on date of announcement – 2025',
                          figw=10.8, dpi=150)

# ── Replace image in slide ────────────────────────────────────────────────────
def replace_slide_image(prs, slide_index, img_buf, _unused):
    """Replace the first pic element image with a new PNG."""
    from pptx.parts.image import ImagePart
    from pptx.opc.packuri import PackURI

    slide = prs.slides[slide_index]
    img_buf.seek(0)
    img_bytes = img_buf.read()

    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            blipFill = shape._element.find('.//' + qn('p:blipFill'))
            if blipFill is None: continue
            blip = blipFill.find(qn('a:blip'))
            if blip is None: continue
            rId = blip.get(qn('r:embed'))
            if rId is None: continue

            # Get the image part and overwrite its blob directly
            img_part = slide.part.related_part(rId)
            img_part._blob = img_bytes
            # Update content type via internal attribute
            img_part._content_type = 'image/png'

            print(f'  ✓ Slide {slide_index+1}: image replaced ({len(img_bytes)//1024}KB)')
            return True

    print(f'  ⚠ Slide {slide_index+1}: no pic shape found')
    return False

# ── Update slide 2 XML tables ─────────────────────────────────────────────────
def update_slide2(prs, agencies, group_stats):
    slide = prs.slides[1]
    tables_found = []
    for shape in slide.shapes:
        if shape.has_table:
            tables_found.append(shape.table)
    print(f'  Slide 2: {len(tables_found)} tables found')

    # ── Table 4 (idx 4): Top agencies ──
    if len(tables_found) > 4:
        tbl  = tables_found[4]
        top6 = agencies[:4]
        for i, ag in enumerate(top6):
            if i + 1 < len(tbl.rows):
                row  = tbl.rows[i + 1]
                vals = [str(ag['rank']),
                        ' ' + ag['agency'].title(),
                        str(int(ag['nbb'])),
                        str(int(ag['wins'])),
                        str(int(ag['dep']))]
                for j, cell in enumerate(list(row.cells)[:5]):
                    if j >= len(vals): break
                    tf = cell.text_frame
                    # Clear ALL runs in ALL paragraphs first
                    for para in tf.paragraphs:
                        for run in para.runs:
                            run.text = ''
                    # Write into first paragraph's first run (or add one)
                    if tf.paragraphs:
                        para = tf.paragraphs[0]
                        if para.runs:
                            para.runs[0].text = vals[j]
                        else:
                            run = para.add_run()
                            run.text = vals[j]
        print('  ✓ Top agencies table updated')

    # ── Table 5 (idx 5): Key Takeaways ──
    if len(tables_found) > 5:
        tbl = tables_found[5]
        top_ag = agencies[0]
        ag2    = agencies[1]
        def gnbb(name): return next((a['nbb'] for a in agencies if a['agency']==name), 0)

        takeaways = (
            f"• Publicis Media leads the market with +${group_stats['Publicis Media']['nbb']:.0f}m NBB: "
            f"{top_ag['agency'].title()} is the top individual winner (+${top_ag['nbb']:.0f}m), "
            f"driven by the Nestlé + Wyeth consolidation. Starcom adds +${gnbb('STARCOM'):.0f}m "
            f"and Zenith +${gnbb('ZENITH'):.0f}m.\n"
            f"• {ag2['agency'].title()} posts +${ag2['nbb']:.0f}m NBB built on Telford ($27.5m) "
            f"alongside a strong retained book (KAO, Ocean Park, Asahi).\n"
            f"• Omnicom consolidates for a combined +${group_stats['Omnicom Media']['nbb']:.0f}m NBB. "
            f"Hearts & Science leads with AS Watson ($18m). PHD delivers +${gnbb('PHD'):.0f}m.\n"
            f"• Havas ends the year effectively flat — Wins (${group_stats['Havas Media Network']['wins']:.0f}m) "
            f"and departures (${group_stats['Havas Media Network']['dep']:.0f}m) largely cancel out.\n"
            f"• WPP Media faces a challenging year with a combined NBB of ${group_stats['WPP Media']['nbb']:.0f}m. "
            f"EssenceMediacom delivers +${gnbb('ESSENCEMEDIACOM'):.0f}m on the back of the Webull win ($6.2m)."
        )
        if len(tbl.rows) > 1:
            cell = tbl.rows[1].cells[0]
            tf   = cell.text_frame
            tf.clear()
            p    = tf.add_paragraph()
            run  = p.add_run()
            run.text = takeaways
            run.font.size = Pt(7.5)
        print('  ✓ Key Takeaways updated')

# ── Main ──────────────────────────────────────────────────────────────────────
def main():
    print('Loading data...')
    agencies, group_stats, top_moves, df, market_name = load_stats()
    print(f'  {len(agencies)} agencies, {len(top_moves)} top moves for {market_name}')

    print('Generating table images...')
    imgs = {
        3: make_slide3_img(top_moves),
        4: make_slide4_img(agencies),
        5: make_slide5_img(agencies, group_stats),
        6: make_slide6_chart(agencies),
        7: make_slide7_img(agencies),
    }
    print('  ✓ All images generated')

    print('Opening template...')
    prs = Presentation(TEMPLATE)

    print(f'Replacing "Hong Kong" with "{market_name}"...')
    replace_text_in_pptx(prs, "Hong Kong", market_name)

    print('Updating slide 2 tables...')
    update_slide2(prs, agencies, group_stats)

    # Slides 3 & 4 : design moderne RECMA (tables natives éditables)
    from modern_design import build_slide3_modern, build_slide4_modern
    print('Building modern RECMA design — Slide 3 (TOP moves)...')
    build_slide3_modern(prs.slides[2], top_moves, market=market_name)
    print('  ✓ Slide 3: modern design applied')

    print('Building modern RECMA design — Slide 4 (T21a agencies)...')
    build_slide4_modern(prs.slides[3], agencies, market=market_name)
    print('  ✓ Slide 4: modern design applied')

    print('Replacing slide images (slides 5, 6, 7)...')
    for slide_idx in [4, 5, 6]:
        replaced = replace_slide_image(prs, slide_idx, imgs[slide_idx+1], None)
        if not replaced and slide_idx == 5:
            slide = prs.slides[5]
            imgs[6].seek(0)
            for shape in slide.shapes:
                if shape.shape_type == 3:  # CHART
                    left, top, w, h = shape.left, shape.top, shape.width, shape.height
                    slide.shapes.add_picture(imgs[6], left, top, w, h)
                    print(f'  ✓ Slide 6: chart covered with image')
    
    output_path = OUTPUT.replace('HK', market_name.replace(' ', '_'))
    print(f'Saving to {output_path}...')
    prs.save(output_path)
    size = os.path.getsize(output_path)
    print(f'✅ Done! {size/1024:.0f} KB → {output_path}')

if __name__ == '__main__':
    main()
